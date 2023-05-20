#うごかない
# https://kakakakakku.hatenablog.com/entry/2022/11/24/105451

import collections.abc
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[10]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "こんにちは世界!"
subtitle.text = "python-pptx はここにありました!"

prs.save('test2.pptx')

text_runs = []

num = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            
            for run in paragraph.runs:
                num+=len(run.text)                

print('文字数:',num)
                
            